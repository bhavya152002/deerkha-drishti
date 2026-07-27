"""
Build the Deerkha Drishti technical review deck (.pptx).

Design: 16:9, light ground (projector-safe), forest-green accent carried over
from the deployment guide so the document set reads as one system.
Typography: Segoe UI (present on any Windows machine) + Consolas for anything
that is literally a path, command or identifier.
"""

from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE

# ----------------------------------------------------------------- palette --
INK       = C(0x1B, 0x20, 0x18)   # near-black, green-biased
DEEP      = C(0x14, 0x18, 0x0F)   # section-divider ground
ACCENT    = C(0x1F, 0x7A, 0x66)   # viridian
ACCENT_LO = C(0xD6, 0xE7, 0xE0)
PAPER     = C(0xFF, 0xFF, 0xFF)
PANEL     = C(0xF3, 0xF5, 0xEF)
SUNK      = C(0xE4, 0xE8, 0xDE)
LINE      = C(0xC7, 0xCD, 0xBE)
DIM       = C(0x5A, 0x63, 0x53)
MUTE      = C(0x7C, 0x84, 0x74)
WARN      = C(0x8A, 0x64, 0x12)
WARN_LO   = C(0xF7, 0xEE, 0xD6)
CRIT      = C(0xA8, 0x3C, 0x24)
CRIT_LO   = C(0xF7, 0xE3, 0xDD)
GOOD      = C(0x3F, 0x7A, 0x48)
GOOD_LO   = C(0xE2, 0xEE, 0xE1)

SANS = "Segoe UI"
MONO = "Consolas"

W, H = 13.333, 7.5
MARGIN = 0.72

prs = Presentation()
prs.slide_width = In(W)
prs.slide_height = In(H)
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------------ helpers --
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER
    return s


def tb(s, x, y, w, h, text, size=14, color=INK, bold=False, font=SANS,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15, italic=False):
    box = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return box


def rect(s, x, y, w, h, fill=PANEL, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         adj=0.06):
    sh = s.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    if fill is None:
        sh.fill.background()          # transparent, e.g. a border-only overlay
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = adj
        except Exception:
            pass
    sh.text_frame.word_wrap = True
    return sh


def node(s, x, y, w, h, title, sub=None, fill=PANEL, line=LINE, tsize=12, ssize=9.5,
         tcolor=INK, scolor=MUTE, tfont=SANS, sfont=MONO):
    """A labelled box for diagrams."""
    rect(s, x, y, w, h, fill=fill, line=line)
    if sub:
        tb(s, x + 0.08, y + h / 2 - 0.30, w - 0.16, 0.30, title, size=tsize, bold=True,
           color=tcolor, font=tfont, align=PP_ALIGN.CENTER)
        tb(s, x + 0.08, y + h / 2 + 0.02, w - 0.16, 0.26, sub, size=ssize,
           color=scolor, font=sfont, align=PP_ALIGN.CENTER)
    else:
        tb(s, x + 0.08, y, w - 0.16, h, title, size=tsize, bold=True, color=tcolor,
           font=tfont, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def arrow(s, x1, y1, x2, y2, color=ACCENT, lw=1.6, dashed=False, head=True):
    ln = s.shapes.add_connector(1, In(x1), In(y1), In(x2), In(y2))  # 1 = straight
    ln.line.color.rgb = color
    ln.line.width = Pt(lw)
    if dashed:
        ln.line.dash_style = MSO_LINE.DASH
    if head:
        ln.line._get_or_add_ln().append(_arrowhead())
    return ln


def _arrowhead():
    from pptx.oxml.ns import qn
    from lxml import etree
    el = etree.SubElement(etree.Element("x"), qn("a:tailEnd"))
    el.set("type", "triangle")
    el.set("w", "med")
    el.set("len", "med")
    return el


def label(s, x, y, w, text, size=9, color=MUTE, font=MONO, align=PP_ALIGN.CENTER):
    return tb(s, x, y, w, 0.24, text, size=size, color=color, font=font, align=align)


# -------------------------------------------------------------- slide kinds --
PAGE = {"n": 0}


def chrome(s, section=None):
    """Footer rule + slide number + running section label."""
    PAGE["n"] += 1
    rect(s, MARGIN, H - 0.62, W - 2 * MARGIN, 0.012, fill=LINE, line=None,
         shape=MSO_SHAPE.RECTANGLE)
    if section:
        tb(s, MARGIN, H - 0.52, 6.0, 0.26, section.upper(), size=8.5, color=MUTE,
           font=MONO)
    tb(s, W - MARGIN - 1.2, H - 0.52, 1.2, 0.26, f"{PAGE['n']:02d}", size=8.5,
       color=MUTE, font=MONO, align=PP_ALIGN.RIGHT)


def title_slide():
    s = slide()
    rect(s, 0, 0, W, H, fill=DEEP, line=None, shape=MSO_SHAPE.RECTANGLE)
    rect(s, 0, 0, 0.18, H, fill=ACCENT, line=None, shape=MSO_SHAPE.RECTANGLE)
    tb(s, 1.15, 2.05, 10.5, 0.4, "DEERKHA DRISHTI", size=13, color=C(0x58, 0xBF, 0xA4),
       font=MONO, bold=True)
    tb(s, 1.15, 2.62, 11, 1.5,
       "System Architecture &\nTechnical Review", size=42, bold=True,
       color=C(0xE2, 0xE7, 0xDC), spacing=1.02)
    rect(s, 1.15, 4.42, 2.4, 0.028, fill=ACCENT, line=None, shape=MSO_SHAPE.RECTANGLE)
    tb(s, 1.15, 4.78, 9.6, 1.0,
       "Wildlife detection, deterrence and alerting across distributed forest sites.\n"
       "How every part connects — and what it takes to scale well past the current target.",
       size=14.5, color=C(0xA2, 0xAC, 0x99), spacing=1.35)
    tb(s, 1.15, 6.35, 9, 0.3, "PREPARED FOR TECHNICAL REVIEW", size=9.5,
       color=C(0x78, 0x82, 0x6F), font=MONO)
    PAGE["n"] += 0
    return s


def section_slide(num, title, blurb):
    s = slide()
    rect(s, 0, 0, W, H, fill=DEEP, line=None, shape=MSO_SHAPE.RECTANGLE)
    tb(s, MARGIN + 0.3, 2.6, 2.0, 1.2, num, size=64, bold=True, color=ACCENT)
    tb(s, MARGIN + 0.3, 3.85, 10.5, 0.8, title, size=32, bold=True,
       color=C(0xE2, 0xE7, 0xDC))
    tb(s, MARGIN + 0.3, 4.85, 8.6, 0.8, blurb, size=13.5,
       color=C(0xA2, 0xAC, 0x99), spacing=1.35)
    PAGE["n"] += 1
    tb(s, W - MARGIN - 1.2, H - 0.52, 1.2, 0.26, f"{PAGE['n']:02d}", size=8.5,
       color=C(0x78, 0x82, 0x6F), font=MONO, align=PP_ALIGN.RIGHT)
    return s


def content(title, kicker=None, section=None):
    s = slide()
    tb(s, MARGIN, 0.52, W - 2 * MARGIN, 0.6, title, size=27, bold=True, color=INK)
    y = 1.20
    if kicker:
        tb(s, MARGIN, 1.16, W - 2 * MARGIN - 0.4, 0.5, kicker, size=13, color=DIM,
           spacing=1.3)
        y = 1.78
    rect(s, MARGIN, y, 1.4, 0.028, fill=ACCENT, line=None, shape=MSO_SHAPE.RECTANGLE)
    chrome(s, section)
    return s, y + 0.30


def bullets(s, x, y, w, items, size=12.5, gap=0.10, bullet_color=ACCENT):
    """items: list of (heading, body) or plain strings."""
    cur = y
    for it in items:
        if isinstance(it, tuple):
            head, body = it
        else:
            head, body = it, None
        rect(s, x, cur + 0.09, 0.10, 0.10, fill=bullet_color, line=None,
             shape=MSO_SHAPE.OVAL)
        hbox = tb(s, x + 0.28, cur, w - 0.28, 0.28, head, size=size, bold=True, color=INK)
        cur += 0.30
        if body:
            bbox = tb(s, x + 0.28, cur, w - 0.28, 0.30, body, size=size - 1.0,
                      color=DIM, spacing=1.28)
            cur += 0.26 * (1 + len(body) // 95)
        cur += gap
    return cur


def table(s, x, y, w, rows, colw, header=True, size=10.5, rowh=0.34, headh=0.32):
    """rows: list of lists of strings. colw: list of fractions summing to 1."""
    widths = [w * f for f in colw]
    cur = y
    for ri, row in enumerate(rows):
        h = headh if (header and ri == 0) else rowh
        cx = x
        if header and ri == 0:
            rect(s, x, cur, w, h, fill=SUNK, line=None, shape=MSO_SHAPE.RECTANGLE)
        elif ri % 2 == 0:
            rect(s, x, cur, w, h, fill=PANEL, line=None, shape=MSO_SHAPE.RECTANGLE)
        for ci, cell in enumerate(row):
            is_head = header and ri == 0
            txt = cell
            fnt = MONO if (txt.startswith("`") and txt.endswith("`")) else SANS
            if fnt == MONO:
                txt = txt[1:-1]
            tb(s, cx + 0.12, cur + (h - 0.22) / 2 - 0.02, widths[ci] - 0.2, 0.24, txt,
               size=8.5 if is_head else size,
               bold=is_head, color=MUTE if is_head else INK, font=MONO if is_head else fnt,
               spacing=1.0)
            cx += widths[ci]
        cur += h
    # border-only overlay so the banded rows read as one table
    rect(s, x, y, w, cur - y, fill=None, line=LINE, shape=MSO_SHAPE.RECTANGLE)
    return cur


def callout(s, x, y, w, kind, lbl, text, size=11.5):
    fills = {"note": (ACCENT_LO, ACCENT), "warn": (WARN_LO, WARN),
             "crit": (CRIT_LO, CRIT), "good": (GOOD_LO, GOOD)}
    bg, fg = fills[kind]
    # Height must come from the ACTUAL width -- a fixed chars-per-line guess
    # silently overflows narrow sidebar callouts.
    import math as _m
    usable_pt = (w - 0.44) * 72
    cpl = max(8, int(usable_pt / (size * 0.52)))
    nlines = max(1, _m.ceil(len(text) / cpl))
    h = 0.50 + 0.215 * nlines
    rect(s, x, y, w, h, fill=bg, line=None, shape=MSO_SHAPE.RECTANGLE)
    rect(s, x, y, 0.045, h, fill=fg, line=None, shape=MSO_SHAPE.RECTANGLE)
    tb(s, x + 0.22, y + 0.11, w - 0.4, 0.22, lbl.upper(), size=8.5, bold=True,
       color=fg, font=MONO)
    tb(s, x + 0.22, y + 0.36, w - 0.44, h - 0.45, text, size=size, color=INK, spacing=1.25)
    return y + h + 0.16


def metric(s, x, y, w, big, lbl, sub=None, color=ACCENT):
    tb(s, x, y, w, 0.62, big, size=30, bold=True, color=color)
    tb(s, x, y + 0.60, w, 0.26, lbl, size=10.5, bold=True, color=INK)
    if sub:
        tb(s, x, y + 0.84, w, 0.5, sub, size=9.5, color=MUTE, spacing=1.22)


# =============================================================== THE DECK ====
title_slide()

# ---- agenda
s, y = content("What this covers", section="Contents")
left = [
    ("A · The system", "What it does, and the three tiers it is built from."),
    ("B · Connection logic", "Every data path: config down, telemetry up, alerts up, video on demand."),
    ("C · Edge internals", "Camera to alert, inside one Jetson."),
    ("D · Server internals", "Request paths, data model, caching."),
]
right = [
    ("E · Deployment & operations", "One artifact, many boxes. Canary, health gate, rollback."),
    ("F · Capacity & scaling", "Measured limits now, and the path well beyond the current target."),
    ("G · Status", "What is built, what is pending, what is deliberately deferred."),
]
bullets(s, MARGIN, y + 0.05, 5.6, left)
bullets(s, MARGIN + 6.2, y + 0.05, 5.6, right)
callout(s, MARGIN, 5.95, W - 2 * MARGIN, "note", "Framing",
        "Figures in this deck are either read directly from the code or derived from it. "
        "Where a number is modelled rather than measured, it is labelled as such.")

# =========================================================== A · THE SYSTEM ==
section_slide("A", "The system",
              "What it does, and the three tiers it is built from.")

s, y = content("The problem, and what the system does",
               "Wildlife enters farmland and forest-edge settlements at night. Detection has to happen "
               "where the cameras are — the connectivity is not good enough to send video anywhere.",
               section="A · The system")
b = bullets(s, MARGIN, y, 5.7, [
    ("Detect", "RTSP cameras watched continuously. Motion gates a TensorRT object detector; a second "
               "vision-language model confirms the species."),
    ("Deter", "Confirmed species trigger relay-driven lights and audio playback, per species."),
    ("Alert", "Telegram message with the annotated frame, plus a durable row in the detection history."),
    ("Record", "Continuous motion video and stills kept locally on a retention window."),
])
node(s, MARGIN + 6.4, y + 0.05, 5.5, 1.05, "All inference happens at the edge",
     "no video leaves the site unless asked for", fill=ACCENT_LO, line=ACCENT, tsize=13)
bullets(s, MARGIN + 6.4, y + 1.35, 5.5, [
    ("Why that matters", "Field links are metered 4G. Streaming raw video to a central detector would be "
                         "both unaffordable and unavailable when the link drops."),
    ("Consequence", "Each site keeps working through a total network outage. The server is for "
                    "configuration, visibility and history — never for the detection path."),
])

s, y = content("Three tiers", section="A · The system")
# tier boxes
tw, th = 3.55, 2.5
tx = [MARGIN, MARGIN + 4.15, MARGIN + 8.30]
tops = y + 0.35
node(s, tx[0], tops, tw, 0.55, "EDGE", "Jetson Orin Nano, one per site",
     fill=ACCENT, line=ACCENT, tcolor=PAPER, scolor=C(0xD6, 0xE7, 0xE0), tsize=13)
node(s, tx[1], tops, tw, 0.55, "SERVER", "single VPS, Mumbai",
     fill=ACCENT, line=ACCENT, tcolor=PAPER, scolor=C(0xD6, 0xE7, 0xE0), tsize=13)
node(s, tx[2], tops, tw, 0.55, "CONSUMERS", "browsers and Telegram",
     fill=ACCENT, line=ACCENT, tcolor=PAPER, scolor=C(0xD6, 0xE7, 0xE0), tsize=13)

body = [
    ["Camera capture and decode",
     "Configuration authority", "Admin dashboard"],
    ["Motion, detection, species confirm",
     "Device telemetry ingest", "Live camera view"],
    ["Deterrence control",
     "Detection history", "Detection history"],
    ["Local recording and retention",
     "Fleet deployment target", "Telegram alerts"],
    ["Alert dispatch", "Backups", "Client dashboard (planned)"],
]
for col in range(3):
    rect(s, tx[col], tops + 0.62, tw, th - 0.62, fill=PANEL, line=LINE)
for r, row in enumerate(body):
    for col in range(3):
        tb(s, tx[col] + 0.22, tops + 0.80 + r * 0.36, tw - 0.4, 0.3, "· " + row[col],
           size=11, color=INK if r == 0 else DIM)

callout(s, MARGIN, tops + th + 0.30, W - 2 * MARGIN, "note", "The load-bearing property",
        "The edge tier is autonomous. It boots and detects on a cached configuration, so neither the "
        "server nor the internet is in the detection path. Everything else in this deck follows from that choice.")

s, y = content("Physical topology", "Where the boxes actually are, and what carries traffic between them.",
               section="A · The system")
# left: site
rect(s, MARGIN, y + 0.15, 4.3, 3.1, fill=PANEL, line=LINE)
tb(s, MARGIN + 0.2, y + 0.28, 3.9, 0.3, "FOREST SITE  ×7", size=9.5, bold=True,
   color=MUTE, font=MONO)
for i in range(3):
    node(s, MARGIN + 0.25, y + 0.72 + i * 0.52, 1.5, 0.42, f"Camera {i+1}", fill=PAPER, tsize=10)
tb(s, MARGIN + 0.30, y + 2.34, 1.5, 0.24, "… up to 8", size=9.5, color=MUTE, font=MONO)
node(s, MARGIN + 2.35, y + 1.05, 1.75, 0.95, "Jetson", "Orin Nano", fill=ACCENT_LO, line=ACCENT)
for i in range(3):
    arrow(s, MARGIN + 1.78, y + 0.93 + i * 0.52, MARGIN + 2.32, y + 1.52, color=MUTE, lw=1.1)
label(s, MARGIN + 0.9, y + 2.62, 3.4, "RTSP over the site's own LAN", color=MUTE)

# middle: tailnet
node(s, MARGIN + 5.05, y + 1.05, 2.3, 0.95, "Tailscale", "WireGuard mesh",
     fill=PANEL, line=ACCENT)
arrow(s, MARGIN + 4.15, y + 1.52, MARGIN + 5.02, y + 1.52)
label(s, MARGIN + 4.0, y + 1.72, 1.2, "4G / LTE", color=MUTE)

# right: vps
rect(s, MARGIN + 8.05, y + 0.15, 3.9, 3.1, fill=PANEL, line=LINE)
tb(s, MARGIN + 8.25, y + 0.28, 3.5, 0.3, "VPS · MUMBAI", size=9.5, bold=True,
   color=MUTE, font=MONO)
node(s, MARGIN + 8.30, y + 0.70, 3.4, 0.5, "nginx + TLS", fill=PAPER, tsize=10.5)
node(s, MARGIN + 8.30, y + 1.30, 3.4, 0.5, "FastAPI application", fill=ACCENT_LO,
     line=ACCENT, tsize=10.5)
node(s, MARGIN + 8.30, y + 1.90, 3.4, 0.5, "PostgreSQL", fill=PAPER, tsize=10.5)
node(s, MARGIN + 8.30, y + 2.50, 3.4, 0.5, "backups + retention timers", fill=PAPER, tsize=10.5)
arrow(s, MARGIN + 7.38, y + 1.52, MARGIN + 8.02, y + 1.52)

arrow(s, MARGIN + 10.0, y + 3.30, MARGIN + 10.0, y + 3.70, color=MUTE)
node(s, MARGIN + 8.30, y + 3.72, 3.4, 0.45, "browsers  ·  HTTPS", fill=PAPER, tsize=10.5)

callout(s, MARGIN, y + 3.60, 7.6, "warn", "One consequence to hold onto",
        "Cameras are reachable only from their own Jetson. Nothing in the cloud can talk to a camera — "
        "which is why live view is a proxied pull rather than a stream the site pushes.")

# ==================================================== B · CONNECTION LOGIC ===
section_slide("B", "Connection logic",
              "Every data path in the system: what initiates it, what carries it, and how often.")

s, y = content("Four data flows — and every one is edge-initiated",
               section="B · Connection logic")
flows = [
    ("1", "Configuration  ·  server → edge", "Edge polls every 30 s. Cheap version check first; the full "
     "blob only when the hash changed.", "PULL", ACCENT),
    ("2", "Telemetry  ·  edge → server", "Heartbeat every 60 s: cameras alive, config version, disk free, "
     "release tag, error counters.", "PUSH", ACCENT),
    ("3", "Detections  ·  edge → database", "One row per confirmed alert, written directly. Telegram sent "
     "in parallel from the edge.", "PUSH", ACCENT),
    ("4", "Live video  ·  browser → server → edge", "Only while somebody is watching. Server proxies; the "
     "browser never reaches a Jetson.", "PULL", WARN),
]
cy = y + 0.05
for num, head, body, kind, col in flows:
    rect(s, MARGIN, cy, W - 2 * MARGIN, 0.92, fill=PANEL, line=LINE)
    rect(s, MARGIN, cy, 0.05, 0.92, fill=col, line=None, shape=MSO_SHAPE.RECTANGLE)
    tb(s, MARGIN + 0.28, cy + 0.20, 0.4, 0.4, num, size=17, bold=True, color=col, font=MONO)
    tb(s, MARGIN + 0.85, cy + 0.15, 6.4, 0.3, head, size=13.5, bold=True, color=INK)
    tb(s, MARGIN + 0.85, cy + 0.48, 8.4, 0.34, body, size=11, color=DIM, spacing=1.2)
    rect(s, W - MARGIN - 1.15, cy + 0.28, 0.95, 0.36, fill=PAPER, line=col)
    tb(s, W - MARGIN - 1.15, cy + 0.35, 0.95, 0.24, kind, size=9, bold=True, color=col,
       font=MONO, align=PP_ALIGN.CENTER)
    cy += 1.04

callout(s, MARGIN, cy + 0.02, W - 2 * MARGIN, "good", "Why this matters for security",
        "No inbound connection to a field site is ever required. Every site sits behind whatever NAT its "
        "4G carrier provides, with no port forwarding and no public address.")

s, y = content("How a setting reaches a camera",
               "Worked example: an operator changes a detection threshold in the dashboard.",
               section="B · Connection logic")
steps = [
    ("Operator saves", "Dashboard writes the row and calls one shared function."),
    ("Version recomputed", "The whole device configuration is reassembled and hashed. That SHA-256 is the "
     "config version, stored on the device row."),
    ("Cache warmed", "The freshly assembled blob is put in an in-process cache keyed by that hash."),
    ("Edge polls", "Within 30 s the device asks for its version — two database round-trips, no assembly."),
    ("Hash differs", "The device fetches the full blob. A cache hit costs no database work at all."),
    ("Applied live", "Immutable snapshot swapped in. Most settings take effect on the next frame."),
    ("Acknowledged", "The device reports the version it is running, which the dashboard displays."),
]
cy = y
for i, (head, body) in enumerate(steps):
    rect(s, MARGIN + 0.02, cy + 0.06, 0.34, 0.34, fill=ACCENT, line=None, shape=MSO_SHAPE.OVAL)
    tb(s, MARGIN + 0.02, cy + 0.12, 0.34, 0.24, str(i + 1), size=10.5, bold=True,
       color=PAPER, font=MONO, align=PP_ALIGN.CENTER)
    tb(s, MARGIN + 0.52, cy + 0.02, 2.9, 0.3, head, size=12.5, bold=True, color=INK)
    tb(s, MARGIN + 3.5, cy + 0.02, 6.4, 0.5, body, size=11, color=DIM, spacing=1.2)
    if i < len(steps) - 1:
        arrow(s, MARGIN + 0.19, cy + 0.42, MARGIN + 0.19, cy + 0.60, color=LINE, lw=1.2, head=False)
    cy += 0.66

callout(s, MARGIN + 10.15, y + 0.1, 2.05, "note", "Restart?",
        "Only detection resolution and model paths need a restart. Everything else — including "
        "adding or removing a camera — applies live.", size=10.5)

s, y = content("Detection path — pixels to alert",
               "Inside one Jetson, for one camera. Every stage exists to avoid doing the next one.",
               section="B · Connection logic")
stages = [
    ("RTSP", "H.265", PANEL),
    ("Hardware\ndecode", "NVDEC", PANEL),
    ("Motion\ngate", "640×360", ACCENT_LO),
    ("Object\ndetector", "RF-DETR TRT", ACCENT_LO),
    ("Species\nconfirm", "MobileCLIP", ACCENT_LO),
    ("Confirm\nwindow", "N frames", PANEL),
    ("Act", "alert + deter", GOOD_LO),
]
bw, bh = 1.48, 1.05
bx = MARGIN
for i, (t, sub, fill) in enumerate(stages):
    node(s, bx, y + 0.35, bw, bh, t, sub, fill=fill,
         line=ACCENT if fill == ACCENT_LO else (GOOD if fill == GOOD_LO else LINE), tsize=11)
    if i < len(stages) - 1:
        arrow(s, bx + bw + 0.02, y + 0.88, bx + bw + 0.22, y + 0.88, lw=1.4)
    bx += bw + 0.24

notes = [
    ("Motion gates everything downstream", "An idle camera runs no inference at all. This is what makes "
     "eight cameras per box possible on a 1024-core GPU."),
    ("Inference is paced, not free-running", "Five inferences per second per camera. Uncapped, the loop "
     "re-ran on the same frame at up to 100 Hz, burning GPU to recompute an identical answer."),
    ("Confirmation is temporal", "A species must persist across a window of frames before it counts. "
     "This is what removes single-frame false positives."),
    ("Results expire", "A published detection is only trusted for ~1.5 s. Without that bound, a camera "
     "keeps acting on a frozen frame after motion ends — which is a real failure this system had."),
]
bullets(s, MARGIN, y + 1.72, 5.6, notes[:2])
bullets(s, MARGIN + 6.2, y + 1.72, 5.6, notes[2:])

s, y = content("Live video — the one path that costs real money",
               section="B · Connection logic")
chain = [("Browser", "HTTPS"), ("nginx", "TLS"), ("FastAPI", "proxy + cache"),
         ("Tailscale", "WireGuard"), ("Jetson", "in-memory frame")]
bx = MARGIN
for i, (t, sub) in enumerate(chain):
    node(s, bx, y + 0.25, 2.15, 0.85, t, sub, fill=PANEL if i not in (2,) else ACCENT_LO,
         line=ACCENT if i == 2 else LINE, tsize=12)
    if i < len(chain) - 1:
        arrow(s, bx + 2.17, y + 0.67, bx + 2.42, y + 0.67, lw=1.4)
    bx += 2.42

cy = y + 1.35
rows = [
    ["Control", "Effect"],
    ["Tile size 480 px at quality 55", "The grid renders ~320 px tiles. Serving 1280 px shipped 16× the pixels displayed."],
    ["Refresh every 5 s, not 1 s", "Five-fold reduction. Sub-second refresh was removed entirely."],
    ["2-second server-side cache with request coalescing", "Cost scales with cameras, not with viewers. Ten people watching one camera cost one fetch from the field."],
    ["Polling stops when the tab is hidden", "Removes the failure that actually happens: a dashboard left open overnight."],
    ["Concurrent video streams capped", "Continuous MJPEG is the expensive mode; it is bounded and admin-only."],
]
table(s, MARGIN, cy, W - 2 * MARGIN, rows, [0.34, 0.66], size=10.5, rowh=0.44)
callout(s, MARGIN, cy + 2.55, W - 2 * MARGIN, "warn", "Why it was worth this much attention",
        "Before these controls an eight-camera grid ran roughly 76 GB per day of site uplink. On a metered "
        "4G link that also starves the config poll and the heartbeat — so the device shows OFFLINE precisely "
        "because somebody is watching it.")

s, y = content("Security boundaries", "What is reachable from where.", section="B · Connection logic")
rows = [
    ["Surface", "Reachable from", "Protected by"],
    ["Admin dashboard", "public internet, HTTPS", "session cookie, `Secure` + signed; per-IP login rate limit"],
    ["Client dashboard (planned)", "public internet, HTTPS", "separate hostname; admin routes hard-blocked at nginx"],
    ["Device API — config, heartbeat", "tailnet only", "per-device bearer token, scoped to that device id"],
    ["Live frame server on the Jetson", "tailnet only", "device token; the browser never sees it"],
    ["PostgreSQL", "loopback + tailnet", "app role; devices get a separate INSERT-only role"],
    ["Application port 8000", "not public", "firewall default-deny; verified by an explicit test"],
    ["Cameras", "site LAN only", "unreachable from the internet by construction"],
]
table(s, MARGIN, y, W - 2 * MARGIN, rows, [0.26, 0.24, 0.50], size=10.5, rowh=0.42)
callout(s, MARGIN, y + 3.35, W - 2 * MARGIN, "good", "Two independent locks on the client surface",
        "Role checks in the application are the real control. The customer-facing hostname additionally "
        "refuses admin and device routes at the proxy, so a future mistake in one layer is not sufficient on its own.")

# ==================================================== C · EDGE INTERNALS =====
section_slide("C", "Edge internals",
              "What runs on one Jetson, and why it is arranged that way.")

s, y = content("Threads and ownership", "Per camera, plus the process-wide workers.",
               section="C · Edge internals")
rect(s, MARGIN, y + 0.1, 6.0, 2.65, fill=PANEL, line=LINE)
tb(s, MARGIN + 0.22, y + 0.24, 5.6, 0.3, "PER CAMERA  ×N", size=9.5, bold=True, color=MUTE, font=MONO)
per = [("Capture thread", "pulls decoded frames from GStreamer"),
       ("Motion + record thread", "motion detection, recording, alert dispatch"),
       ("Inference thread", "paced object detection and species confirm")]
for i, (t, sub) in enumerate(per):
    node(s, MARGIN + 0.25, y + 0.62 + i * 0.66, 5.5, 0.56, t, sub, fill=PAPER, tsize=11.5, ssize=9.5)

rect(s, MARGIN + 6.4, y + 0.1, 5.5, 2.65, fill=PANEL, line=LINE)
tb(s, MARGIN + 6.62, y + 0.24, 5.1, 0.3, "PROCESS-WIDE", size=9.5, bold=True, color=MUTE, font=MONO)
glob = [("Config poller", "30 s, with backoff"), ("Heartbeat", "60 s"),
        ("Camera reconciler", "applies roster changes live"),
        ("Storage cleanup", "every 30 min"), ("Live frame server", "serves already-decoded frames")]
for i, (t, sub) in enumerate(glob):
    node(s, MARGIN + 6.65, y + 0.62 + i * 0.42, 5.2, 0.36, f"{t} — {sub}", fill=PAPER, tsize=10.5)

cy = callout(s, MARGIN, y + 2.95, W - 2 * MARGIN, "note", "Deliberate: one process, not one per camera",
             "Per-camera processes would be the obvious way to escape Python's global lock — but each would "
             "carry its own copy of the inference engines and CUDA context, roughly 2.5 GB. On an 8 GB board "
             "that is an immediate out-of-memory failure.")
callout(s, MARGIN, cy, W - 2 * MARGIN, "good", "Live frame server reuses decoded frames",
        "The live view never opens a second decode of the camera. It publishes frames the detection pipeline "
        "has already produced, so viewing costs encode time only.")

s, y = content("Recording, storage and retention", section="C · Edge internals")
metric(s, MARGIN, y + 0.1, 2.6, "~21 GB", "per camera per day",
       "modelled: 1080p H.265 at 3 Mbps, 60% motion duty")
metric(s, MARGIN + 3.0, y + 0.1, 2.6, "~168 GB", "per box per day", "at 8 cameras")
metric(s, MARGIN + 6.0, y + 0.1, 2.6, "465 GB", "NVMe fitted per box", "the actual field hardware")
metric(s, MARGIN + 9.0, y + 0.1, 2.9, "~2.4 days", "retention that buys", "at 8 cameras", color=WARN)
cy = y + 1.55
b1 = bullets(s, MARGIN, cy, 5.6, [
    ("Two rules, every pass", "Delete folders older than the window; then, while free space is below the "
     "floor, delete the oldest remaining."),
    ("Night-window aware", "Between midnight and dawn the pipeline still writes into the previous day's "
     "folder, so cleanup mirrors that arithmetic exactly rather than trusting the calendar."),
])
b2 = bullets(s, MARGIN + 6.2, cy, 5.6, [
    ("Runs in-process, not on a timer", "An external job cannot know which folder is being written to "
     "right now. Deleting a live folder under space pressure at 3am is real data loss."),
    ("Evidence is prioritised", "Under pressure, bulk motion video is reclaimed before confirmed-detection "
     "footage and stills."),
])
callout(s, MARGIN, max(b1, b2) + 0.10, W - 2 * MARGIN, "warn",
        "Retention is set by camera bitrate, not by disk size",
        "At a 3 Mbps main stream the fitted 465 GB NVMe holds about 2.4 days across 8 cameras, and closer "
        "to 1.5 days on a night when motion runs continuously. Reducing the camera's own encoder to 2 Mbps "
        "restores a full 3-day window at no cost in compute and little visible quality loss. That is a "
        "camera setting, not a code change — which is why retention is quoted as a target here, not a "
        "guarantee, until a week of field measurement replaces the model.")

# ================================================== D · SERVER INTERNALS =====
section_slide("D", "Server internals",
              "Request paths, the data model, and the caching that keeps it cheap.")

s, y = content("What the server actually does per request", section="D · Server internals")
rows = [
    ["Endpoint", "Frequency", "Cost", "Notes"],
    ["Config version check", "1 per device per 30 s", "2 queries", "Serves a stored hash; no assembly"],
    ["Full config fetch", "only when changed", "cache hit → 0 queries", "Content-addressed by hash"],
    ["Heartbeat ingest", "1 per device per 60 s", "1 insert", "Runs off the event loop"],
    ["Live snapshot", "1 per camera per 5 s per viewer", "0 queries", "Device lookup and image both cached"],
    ["Detection history", "on demand", "1 indexed query", "Index added; previously a full table scan"],
]
table(s, MARGIN, y, W - 2 * MARGIN, rows, [0.26, 0.24, 0.20, 0.30], size=10.5, rowh=0.42)
cy = y + 2.55
metric(s, MARGIN, cy, 3.0, "~0.35 req/s", "total, at 7 devices",
       "The server is not the constraint and will not be for a long time.")
callout(s, MARGIN + 3.6, cy - 0.05, 8.3, "note", "Configuration assembly used to dominate",
        "Building one device's configuration is eight sequential queries plus a hash. Served from a hosted "
        "database ~130 ms away that was roughly a second per assembly, and it ran on every poll. It is now "
        "a stored value plus a cache, against a database on the same host.")

s, y = content("The cache that makes it scale", "A small design decision with disproportionate consequences.",
               section="D · Server internals")
node(s, MARGIN + 1.2, y + 0.35, 4.2, 0.95, "Cache key", "(device id, config version)",
     fill=ACCENT_LO, line=ACCENT, tsize=14)
arrow(s, MARGIN + 5.45, y + 0.82, MARGIN + 6.15, y + 0.82)
node(s, MARGIN + 6.2, y + 0.35, 4.6, 0.95, "…and the version IS the content hash",
     "so the key contains the value's identity", fill=PANEL, line=LINE, tsize=12.5)
bullets(s, MARGIN, y + 1.65, 5.6, [
    ("It cannot serve stale data", "If the configuration changes the hash changes, so the old key is "
     "simply never requested again."),
    ("No invalidation protocol", "Nothing to get wrong, no TTL to tune, no cross-process messaging."),
])
bullets(s, MARGIN + 6.2, y + 1.65, 5.6, [
    ("Safe across multiple workers", "A second process with a cold cache does more work — never wrong "
     "work. This is what makes horizontal scaling free later."),
    ("Bounded", "One entry per device. Writing a new version evicts the old one."),
])
callout(s, MARGIN, y + 3.35, W - 2 * MARGIN, "good", "Scaling relevance",
        "This is the single reason the server can be scaled out horizontally without designing a cache "
        "coherence scheme first. It is cited again in section F.")

s, y = content("Data model", "Ten configuration tables plus the detection history. One partition key throughout.",
               section="D · Server internals")
groups = [
    ("Identity", ["devices", "cameras"], ACCENT_LO),
    ("Detection tuning", ["class settings", "CLIP prompts", "CLIP distractors", "global settings"], PANEL),
    ("Deterrence", ["deterrence global", "deterrence targets"], PANEL),
    ("Operational", ["device status", "audit log", "detections"], PANEL),
]
gx = MARGIN
for gname, items, fill in groups:
    gw = 2.9
    rect(s, gx, y + 0.1, gw, 2.3, fill=PAPER, line=LINE)
    tb(s, gx + 0.18, y + 0.24, gw - 0.36, 0.3, gname.upper(), size=9, bold=True, color=MUTE, font=MONO)
    for i, it in enumerate(items):
        node(s, gx + 0.2, y + 0.62 + i * 0.42, gw - 0.4, 0.34, it, fill=fill,
             line=ACCENT if fill == ACCENT_LO else LINE, tsize=10.5)
    gx += 3.05
callout(s, MARGIN, y + 2.62, W - 2 * MARGIN, "note", "Device id is the partition key everywhere",
        "Every configuration table is keyed by device. That is why a device's whole configuration assembles "
        "with one predictable set of queries — and why sharding by device is a natural future move rather "
        "than a rewrite.")

# ============================================ E · DEPLOYMENT & OPERATIONS ====
section_slide("E", "Deployment & operations",
              "One code artifact, many boxes — and what happens when something fails.")

s, y = content("One artifact, many boxes", "Every Jetson runs byte-identical code.",
               section="E · Deployment")
rect(s, MARGIN, y + 0.15, 5.6, 2.5, fill=PANEL, line=LINE)
tb(s, MARGIN + 0.22, y + 0.3, 5.2, 0.3, "ON EVERY BOX", size=9.5, bold=True, color=MUTE, font=MONO)
paths = [("releases/<tag>/", "immutable code, one dir per release"),
         ("current →", "symlink; a deploy flips this atomically"),
         ("models/", "inference engines — survive deploys"),
         ("/var/lib/…", "config cache — survives deploys"),
         ("/mnt/data/…", "recordings — never touched")]
for i, (p, d) in enumerate(paths):
    tb(s, MARGIN + 0.28, y + 0.70 + i * 0.38, 2.1, 0.28, p, size=10.5, color=INK, font=MONO)
    tb(s, MARGIN + 2.5, y + 0.70 + i * 0.38, 3.0, 0.28, d, size=10, color=DIM)

rect(s, MARGIN + 6.2, y + 0.15, 5.7, 2.5, fill=ACCENT_LO, line=ACCENT)
tb(s, MARGIN + 6.42, y + 0.3, 5.3, 0.3, "THE ONLY PER-DEVICE FILE", size=9.5, bold=True,
   color=ACCENT, font=MONO)
tb(s, MARGIN + 6.42, y + 0.72, 5.2, 1.0,
   "device id\napi token\nserver address", size=13, color=INK, font=MONO, spacing=1.5)
tb(s, MARGIN + 6.42, y + 1.95, 5.2, 0.6,
   "Everything else — cameras, credentials, regions of interest,\nthresholds, retention — comes from the server.",
   size=10.5, color=DIM, spacing=1.25)
callout(s, MARGIN, y + 2.88, W - 2 * MARGIN, "good", "Why this matters at fleet scale",
        "Two boxes can be diffed and must match exactly. Configuration drift between sites — the classic "
        "way a fleet becomes unmaintainable — is structurally impossible rather than merely discouraged.")

s, y = content("Deploying to the fleet", "One command. It is designed to stop rather than to finish.",
               section="E · Deployment")
seq = [("Preflight", "reachable, disk free, identity present"),
       ("Stage", "ship the tagged tree to a temp dir"),
       ("Swap", "atomic symlink flip, restart"),
       ("Health gate", "service active, not restart-looping, fresh heartbeat"),
       ("Canary passes", "only then continue to the rest")]
bx = MARGIN
for i, (t, sub) in enumerate(seq):
    node(s, bx, y + 0.3, 2.15, 0.95, t, sub, fill=PANEL if i < 3 else ACCENT_LO,
         line=ACCENT if i >= 3 else LINE, tsize=12, ssize=9)
    if i < len(seq) - 1:
        arrow(s, bx + 2.17, y + 0.77, bx + 2.42, y + 0.77, lw=1.4)
    bx += 2.42
arrow(s, MARGIN + 7.3, y + 1.30, MARGIN + 7.3, y + 1.75, color=CRIT, lw=1.6)
node(s, MARGIN + 5.6, y + 1.78, 3.4, 0.5, "fails → roll back, halt rollout",
     fill=CRIT_LO, line=CRIT, tsize=11)
bullets(s, MARGIN, y + 2.55, 5.6, [
    ("Canary by default", "One box first, health-gated. A bad release reaching all seven sites at once is "
     "the failure this exists to prevent — each one is a drive into a forest."),
    ("Sequential, not parallel", "A failure stops the rollout instead of racing ahead of it."),
])
bullets(s, MARGIN + 6.2, y + 2.55, 5.6, [
    ("Rollback is a symlink", "The previous release is still on disk. Reverting is flipping a pointer and "
     "restarting — seconds, not a redeploy."),
    ("Offline boxes are reported", "Push deployment cannot update a powered-off box. It is named in the "
     "summary rather than silently skipped."),
])

s, y = content("Failure modes", "What actually happens when each part fails.", section="E · Operations")
rows = [
    ["What fails", "Immediate effect", "Detection continues?", "Recovery"],
    ["Internet at a site", "No alerts, no dashboard visibility", "Yes — fully", "Automatic on reconnect"],
    ["Server or database", "No config changes, no dashboard", "Yes — cached config", "Automatic; devices back off and resync"],
    ["A Jetson", "That site goes dark", "No, that site only", "Service auto-restarts; unlimited retries"],
    ["One camera", "That view is lost", "Yes — others unaffected", "Pipeline rebuilt automatically"],
    ["A bad release", "Caught on the canary", "Yes — rollout halted", "Automatic rollback"],
    ["Disk fills", "Oldest footage reclaimed", "Yes", "Two-rule cleanup, floor-protected"],
]
table(s, MARGIN, y, W - 2 * MARGIN, rows, [0.22, 0.30, 0.22, 0.26], size=10.5, rowh=0.42)
callout(s, MARGIN, y + 3.35, W - 2 * MARGIN, "good", "The pattern",
        "There is no single failure that stops detection everywhere. The only component whose loss takes a "
        "site off the air is that site's own hardware.")

# ============================================== F · CAPACITY AND SCALING =====
section_slide("F", "Capacity & scaling",
              "Measured limits today, and the engineering path well beyond the current target.")

s, y = content("Where the limits actually are", "Per Jetson. None of these are hardware ceilings.",
               section="F · Capacity")
rows = [
    ["Constraint", "Status", "Headroom"],
    ["Hardware video decode", "~60% utilised at 8 cameras", "Not the limit"],
    ["Inference throughput", "Paced; concurrency work planned", "Roughly 3× available"],
    ["Memory (8 GB shared)", "~4.6 GB at 8 cameras", "Tight but sufficient"],
    ["Video encoding", "The real ceiling today", "Removed by recording the camera stream directly"],
    ["Disk capacity (465 GB)", "~2.4 days at 8 cameras", "3 days at a lower camera bitrate"],
    ["Site uplink", "Bounded by design", "Not the limit"],
]
table(s, MARGIN, y, 7.4, rows, [0.36, 0.34, 0.30], size=10.5, rowh=0.42)
mx = MARGIN + 7.9
tb(s, mx, y - 0.05, 4.0, 0.3, "CAMERAS PER BOX", size=9, bold=True, color=MUTE, font=MONO)
ladder = [("Today", "3", MUTE), ("After current work", "5–6", ACCENT), ("After planned work", "8", GOOD)]
for i, (lbl_, val, col) in enumerate(ladder):
    rect(s, mx, y + 0.35 + i * 0.78, 4.0, 0.66, fill=PANEL, line=LINE)
    tb(s, mx + 0.2, y + 0.52 + i * 0.78, 2.6, 0.3, lbl_, size=11.5, color=INK)
    tb(s, mx + 2.9, y + 0.44 + i * 0.78, 1.0, 0.4, val, size=19, bold=True, color=col,
       align=PP_ALIGN.RIGHT)
callout(s, mx, y + 2.78, 4.0, "note", "Read this way",
        "Every gain listed is software. The board is not the bottleneck.")

s, y = content("Scaling beyond the current target",
               "50 cameras across 7 boxes is the contract. This is what the same architecture does past that.",
               section="F · Scaling")
tiers = [
    ("50 cams\n7 boxes", "Current target", "No changes needed beyond planned work.", GOOD_LO, GOOD),
    ("100 cams\n13 boxes", "No architectural change", "Server load still under 1 req/s. Batch the "
     "deployment; collapse one dashboard query.", GOOD_LO, GOOD),
    ("250 cams\n32 boxes", "Two known changes", "Run multiple server workers — already safe because of "
     "the cache design. Parallelise deployment.", WARN_LO, WARN),
    ("500+ cams\n60+ boxes", "Structural moves available", "Shard by device id — already the partition "
     "key. Read replica. Devices self-update instead of being pushed to.", WARN_LO, WARN),
]
cx = MARGIN
for head, sub, body, fill, line in tiers:
    rect(s, cx, y + 0.1, 2.85, 2.75, fill=fill, line=line)
    tb(s, cx + 0.2, y + 0.28, 2.5, 0.75, head, size=17, bold=True, color=INK, spacing=1.1)
    tb(s, cx + 0.2, y + 1.10, 2.5, 0.3, sub, size=10, bold=True, color=line, font=MONO)
    tb(s, cx + 0.2, y + 1.48, 2.5, 1.2, body, size=10.5, color=DIM, spacing=1.25)
    cx += 3.02
callout(s, MARGIN, y + 3.05, W - 2 * MARGIN, "good", "The point",
        "Nothing on this path is a rewrite. Each step is a known, bounded change — and the two that matter "
        "most (horizontal server scaling and sharding) are already unblocked by decisions taken in the "
        "current design.")

s, y = content("Why it scales — six structural properties", section="F · Scaling")
props = [
    ("Edge autonomy", "Sites do not depend on the server to function, so adding sites adds no coupling."),
    ("Stateless server", "All durable state is in the database or on disk. Additional servers are additive."),
    ("Content-addressed cache", "Correct across processes by construction, so horizontal scaling needs no "
     "coherence design."),
    ("Device-keyed partitioning", "One partition key throughout — the natural sharding boundary already exists."),
    ("Immutable releases", "Deployment cost per box is constant and parallelisable; rollback is a pointer flip."),
    ("Mesh networking", "No VPN concentrator to become a bottleneck. Adding a site is a constant-time operation."),
]
cy = y + 0.05
for i, (h, b) in enumerate(props):
    col = MARGIN if i % 2 == 0 else MARGIN + 6.2
    row = i // 2
    ry = cy + row * 1.35
    rect(s, col, ry, 5.6, 1.15, fill=PANEL, line=LINE)
    rect(s, col, ry, 0.045, 1.15, fill=ACCENT, line=None, shape=MSO_SHAPE.RECTANGLE)
    tb(s, col + 0.28, ry + 0.16, 5.1, 0.3, h, size=13, bold=True, color=INK)
    tb(s, col + 0.28, ry + 0.52, 5.1, 0.6, b, size=10.5, color=DIM, spacing=1.25)

s, y = content("Known limits and honest gaps", "What a reviewer should push on.", section="F · Scaling")
rows = [
    ["Gap", "Impact", "Position"],
    ["Single shared admin credential", "No per-user audit or revocation", "User accounts with roles are the next work item"],
    ["Devices hold a database credential", "A stolen box is a wider exposure", "Reduced to insert-only; routing alerts via the server removes it entirely"],
    ["Alert imagery is not centralised", "History shows text, not photographs", "Designed; ships with the alert-routing change"],
    ["Deployment is push-based", "A powered-off box is not updated", "Reported, not silent. Self-update is the answer past ~20 boxes"],
    ["Single server instance", "One host to lose", "Backed up and restorable; horizontal scaling already unblocked"],
    ["Disk figures are modelled", "Retention window unproven at 8 cameras", "One instrumented week of soak replaces the estimate"],
]
table(s, MARGIN, y, W - 2 * MARGIN, rows, [0.28, 0.30, 0.42], size=10.5, rowh=0.42)
callout(s, MARGIN, y + 3.35, W - 2 * MARGIN, "note", "Stated deliberately",
        "Every item above is known, bounded and sequenced. None of them blocks running the current fleet; "
        "the first two block giving an external customer a login.")

# ============================================================ G · STATUS =====
section_slide("G", "Status",
              "What is built, what is pending, and in what order.")

s, y = content("Built versus pending", section="G · Status")
tb(s, MARGIN, y, 5.6, 0.3, "COMPLETE", size=9.5, bold=True, color=GOOD, font=MONO)
done = ["Edge pipeline correctness and pacing", "Device-agnostic single artifact",
        "Fleet provisioning and deployment tooling", "Canary, health gate, automatic rollback",
        "Server caching and event-loop correctness", "Bandwidth controls on live video",
        "Retention, backups and restore testing", "Route-level authorisation with enforcing tests",
        "Self-hosted database design and installer"]
for i, d in enumerate(done):
    rect(s, MARGIN + 0.02, y + 0.42 + i * 0.42, 0.22, 0.22, fill=GOOD, line=None, shape=MSO_SHAPE.OVAL)
    tb(s, MARGIN + 0.4, y + 0.38 + i * 0.42, 5.2, 0.3, d, size=11.5, color=INK)

tb(s, MARGIN + 6.2, y, 5.6, 0.3, "PENDING", size=9.5, bold=True, color=WARN, font=MONO)
todo = [("Record the camera stream directly", "removes the encoding ceiling — 8 cameras per box"),
        ("Inference concurrency", "roughly 3× throughput"),
        ("Route alerts through the server", "removes the last credential from field boxes"),
        ("User accounts and roles", "prerequisite for any external login"),
        ("Client dashboard", "read-only, scoped per customer"),
        ("Alert photographs", "ships with alert routing")]
for i, (t_, sub) in enumerate(todo):
    rect(s, MARGIN + 6.22, y + 0.42 + i * 0.66, 0.22, 0.22, fill=WARN, line=None, shape=MSO_SHAPE.OVAL)
    tb(s, MARGIN + 6.6, y + 0.38 + i * 0.66, 5.2, 0.3, t_, size=11.5, color=INK)
    tb(s, MARGIN + 6.6, y + 0.66 + i * 0.66, 5.2, 0.3, sub, size=10, color=MUTE)

s, y = content("Sequence", section="G · Status")
phases = [
    ("NOW", "Deploy the current build", "Server on its own host, first Jetson in the field, one week of "
     "instrumented soak to replace modelled disk figures with measured ones.", ACCENT),
    ("NEXT", "Reach 8 cameras per box", "Direct stream recording and inference concurrency. This is what "
     "makes 50 cameras fit in 7 boxes rather than 10.", ACCENT),
    ("THEN", "Close the security gaps", "Alert routing through the server, user accounts with roles. "
     "Both are prerequisites for handing a customer a login.", WARN),
    ("AFTER", "Customer-facing dashboard", "Read-only, scoped to a customer's own sites, with alert "
     "photographs.", WARN),
]
cy = y + 0.05
for tag, head, body, col in phases:
    rect(s, MARGIN, cy, W - 2 * MARGIN, 1.16, fill=PANEL, line=LINE)
    rect(s, MARGIN, cy, 0.05, 1.16, fill=col, line=None, shape=MSO_SHAPE.RECTANGLE)
    tb(s, MARGIN + 0.3, cy + 0.20, 1.3, 0.3, tag, size=11, bold=True, color=col, font=MONO)
    tb(s, MARGIN + 1.8, cy + 0.16, 4.3, 0.35, head, size=14, bold=True, color=INK)
    tb(s, MARGIN + 1.8, cy + 0.58, 9.8, 0.5, body, size=11, color=DIM, spacing=1.25)
    cy += 1.28

# ---- close
s = slide()
rect(s, 0, 0, W, H, fill=DEEP, line=None, shape=MSO_SHAPE.RECTANGLE)
rect(s, 0, 0, 0.18, H, fill=ACCENT, line=None, shape=MSO_SHAPE.RECTANGLE)
tb(s, 1.15, 2.35, 10.6, 0.9, "In one sentence", size=30, bold=True, color=C(0xE2, 0xE7, 0xDC))
tb(s, 1.15, 3.35, 10.4, 1.6,
   "Detection is autonomous at each site, configuration is centralised and versioned,\n"
   "every box runs identical code deployed with a canary and automatic rollback —\n"
   "and the parts that would have to change to go well past 50 cameras are already unblocked.",
   size=16, color=C(0xA2, 0xAC, 0x99), spacing=1.45)
rect(s, 1.15, 5.25, 2.4, 0.028, fill=ACCENT, line=None, shape=MSO_SHAPE.RECTANGLE)
tb(s, 1.15, 5.6, 9, 0.4, "QUESTIONS", size=12, color=C(0x78, 0x82, 0x6F), font=MONO, bold=True)
PAGE["n"] += 1

out = "Deerkha-Drishti-Technical-Review.pptx"
prs.save(out)
print(f"saved {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
